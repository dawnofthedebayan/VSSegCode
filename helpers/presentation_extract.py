import os
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import re
from tqdm import tqdm
def extract_case_number(text):
    match = re.search(r'Case (\d+)', text)
    if match:
        return match.group(1)
    return None

def save_images_from_slide(slide, base_output_dir):
    slide_number = slide.slide_id
    slide_output_dir = base_output_dir
    
    if not os.path.exists(slide_output_dir):
        os.makedirs(slide_output_dir)
    
    image_count = 1

    # Extract and process text from the slide
    text = ""
    for text_shape in slide.shapes:
        if text_shape.has_text_frame:
            text += text_shape.text + "\n"
    
    case_number = extract_case_number(text)
    
    if case_number:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type for images in python-pptx
                image = shape.image
                image_bytes = image.blob
                image_name = f"CASE_{case_number}_IMG_{image_count}"
                
                if "T1" in text:
                    #create file name based on case number and image type

                    image_name = image_name + "_T1.png"
                elif "T2" in text:
                    
                    image_name = image_name  + "_T2.png"

                elif "Flair" in text:
                    image_name = image_name + "_T2.png"
                else:
                    image_name = image_name + "_other.png"
                
                
                
                image_path = os.path.join(base_output_dir, image_name)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
    
                print(f"Saved image {image_path}")
                image_count += 1

def main():
    presentation_path = "/media/debayan/c7b64c90-ca4e-4192-8ed9-8fea1d005196/SchwanommaDS/UKE/"  # Replace with your presentation file path

    ppt_files = []
    #Select all files with extension .pptx and store in array 
    for file in os.listdir(presentation_path):
        if file.endswith(".pptx"):
            ppt_files.append(os.path.join(presentation_path, file))

    print(f"Found {len(ppt_files)} presentation files")    
    base_output_dir = "/media/debayan/c7b64c90-ca4e-4192-8ed9-8fea1d005196/SchwanommaDS/UKE/dataset/"  # Base directory to save organized images

    if not os.path.exists(base_output_dir):
        os.makedirs(base_output_dir)

    for ppt_file in tqdm(ppt_files,total=len(ppt_files)):
        
        
        presentation = Presentation(ppt_file)

        for slide in presentation.slides:
            save_images_from_slide(slide, base_output_dir)
    
    

    print("Image extraction and organization completed!")

if __name__ == "__main__":
    main()